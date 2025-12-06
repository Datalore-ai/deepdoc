import json
import base64
import io
import os
import logging
from pathlib import Path
from typing import List, Optional, Dict, Any

# prefer using fitz (PyMuPDF) directly
import fitz
from PIL import Image
from docx import Document
from pptx import Presentation
import pdfplumber

# Mistral client (lazy import)
try:
    from mistralai import Mistral
except Exception:
    Mistral = None  # library may not be installed

# Create a module logger
logger = logging.getLogger(__name__)
if not logger.handlers:
    handler = logging.StreamHandler()
    formatter = logging.Formatter("%(levelname)s:%(message)s")
    handler.setFormatter(formatter)
    logger.addHandler(handler)

# suppress noisy logs — only warnings/errors show
logger.setLevel("WARNING")

# Constants
SUPPORTED_EXTENSIONS = {
    "pdf", "jpg", "jpeg", "png", "gif", "webp", "bmp",
    "txt", "md", "doc", "docx", "pptx"
}
MAX_BASE64_ENCODE_BYTES = 50 * 1024 * 1024  # 50MB
TEXT_LINES_PER_PAGE = 40
PARAS_PER_PAGE = 20
OCR_MODEL = "mistral-ocr-latest"

# Track skipped reasons
SKIPPED_DETAILS: List[Dict[str, str]] = []


# Mistral client creation
def _create_mistral_client(api_key: Optional[str] = None):
    api_key = api_key or os.getenv("MISTRAL_API_KEY")
    if Mistral is None:
        return None
    if not api_key:
        return None
    try:
        return Mistral(api_key=api_key)
    except Exception:
        return None


def encode_pdf(pdf_bytes: bytes) -> Optional[str]:
    """Encode PDF bytes to base64 with size guard."""
    try:
        if len(pdf_bytes) > MAX_BASE64_ENCODE_BYTES:
            return None
        return base64.b64encode(pdf_bytes).decode("utf-8")
    except Exception:
        return None


def convert_to_pdf(file_bytes: bytes, filename: str) -> Optional[bytes]:
    """Convert supported formats into a PDF, else skip."""
    extension = filename.lower().split('.')[-1]

    if extension not in SUPPORTED_EXTENSIONS:
        SKIPPED_DETAILS.append({"file": filename, "reason": "Unsupported file type"})
        return None

    buffer = io.BytesIO()

    try:
        if extension == "pdf":
            return file_bytes

        if extension in {"jpg", "jpeg", "png", "gif", "webp", "bmp"}:
            img = Image.open(io.BytesIO(file_bytes)).convert("RGB")
            img.save(buffer, format="PDF")
            return buffer.getvalue()

        if extension in {"txt", "md"}:
            pdf = fitz.open()
            lines = file_bytes.decode("utf-8", errors="ignore").splitlines()
            for i in range(0, len(lines), TEXT_LINES_PER_PAGE):
                page = pdf.new_page()
                page.insert_text((72, 72), "\n".join(lines[i:i + TEXT_LINES_PER_PAGE]))
            pdf.save(buffer)
            return buffer.getvalue()

        if extension in {"doc", "docx"}:
            doc = Document(io.BytesIO(file_bytes))
            pdf = fitz.open()
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            for i in range(0, len(paragraphs), PARAS_PER_PAGE):
                page = pdf.new_page()
                page.insert_text((72, 72), "\n".join(paragraphs[i:i + PARAS_PER_PAGE]))
            pdf.save(buffer)
            return buffer.getvalue()

        if extension == "pptx":
            prs = Presentation(io.BytesIO(file_bytes))
            pdf = fitz.open()
            for slide in prs.slides:
                text = "\n".join(
                    shape.text for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text
                )
                page = pdf.new_page()
                page.insert_text((72, 72), text)
            pdf.save(buffer)
            return buffer.getvalue()

        SKIPPED_DETAILS.append({"file": filename, "reason": "Unsupported file type"})
        return None

    except Exception as e:
        SKIPPED_DETAILS.append({"file": filename, "reason": f"Conversion failed: {e}"})
        return None


def process_page(idx, ocr_response=None):
    try:
        if ocr_response and hasattr(ocr_response, "pages") and idx < len(ocr_response.pages):
            page = ocr_response.pages[idx]
            return getattr(page, "markdown", None) or getattr(page, "text", None) or ""
        return ""
    except Exception as e:
        return f"Error processing page {idx + 1}: {e}"


def extract_text_from_pdf(pdf_bytes: bytes, advanced=True, mistral_client=None):
    extracted = []

    # basic extraction
    if not advanced:
        try:
            with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
                return [page.get_text() for page in doc]
        except Exception as e:
            return [f"Error during simple extraction: {e}"]

    # count pages
    try:
        with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
            total_pages = len(pdf.pages)
    except Exception:
        try:
            with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
                total_pages = len(doc)
        except Exception as e:
            return [f"Error reading PDF: {e}"]

    # check OCR availability
    encoded = encode_pdf(pdf_bytes)
    client = mistral_client or _create_mistral_client()

    if not encoded or not client:
        return extract_text_from_pdf(pdf_bytes, advanced=False)

    # OCR
    try:
        response = client.ocr.process(
            model=OCR_MODEL,
            document={"type": "document_url", "document_url": f"data:application/pdf;base64,{encoded}"},
            include_image_base64=True
        )
    except Exception as e:
        return [f"Error during OCR: {e}"]

    for idx in range(total_pages):
        extracted.append(process_page(idx, response))

    return extracted


def pretty_print_summary(processed: int, skipped: int, errors: int, skipped_details):
    """Very short summary shown once. Shows OCR status only (Enabled/Disabled)."""
    ocr_enabled = bool(os.getenv("MISTRAL_API_KEY"))
    print("\nProcessed File Summary")
    print(f"Processed: {processed}  |  Skipped: {skipped}  |  Errors: {errors}")
    print(f"OCR: {'Enabled' if ocr_enabled else 'Disabled'}")
    if skipped_details:
        print("\nSkipped files:")
        for d in skipped_details:
            print(f" - {d['file']}: {d['reason']}")
    print("\n")


def create_chunks(directory_path: str):
    """Convert files to PDF, extract pages, return chunks."""
    p = Path(directory_path)
    if not p.exists() or not p.is_dir():
        raise ValueError(f"Directory not found: {directory_path}")

    file_paths = [
        os.path.join(directory_path, f)
        for f in os.listdir(directory_path)
        if os.path.isfile(os.path.join(directory_path, f))
    ]

    Chunks = []
    errors = []
    skipped_files = []
    processed_count = 0
    SKIPPED_DETAILS.clear()

    for file_path in file_paths:
        filename = os.path.basename(file_path)
        extension = filename.split(".")[-1].lower()

        try:
            with open(file_path, "rb") as f:
                file_bytes = f.read()
        except Exception as e:
            errors.append({"file": filename, "reason": f"Read error: {e}"})
            continue

        converted = convert_to_pdf(file_bytes, filename)
        if converted is None:
            skipped_files.append(filename)
            continue

        processed_count += 1

        try:
            if extension in {"txt", "md"}:
                pages = extract_text_from_pdf(converted, advanced=False)
            else:
                pages = extract_text_from_pdf(converted, advanced=True)
        except Exception as e:
            errors.append({"file": filename, "reason": f"Extraction error: {e}"})
            pages = [f"Error extracting: {e}"]

        for i, content in enumerate(pages, 1):
            Chunks.append({
                "filename": filename,
                "page_number": i,
                "page_content": content
            })

    # Final summary
    pretty_print_summary(
        processed_count,
        len(skipped_files),
        len(errors),
        SKIPPED_DETAILS
    )

    return Chunks
