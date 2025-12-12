import io, subprocess, tempfile
from pathlib import Path
from typing import Optional
from PIL import Image
from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF

TEXT_LINES_PER_PAGE = 40
PARAS_PER_PAGE = 20

 # Convert .doc to .docx using LibreOffice (Linux only).
def convert_doc_to_docx_linux(file_bytes: bytes) -> Optional[bytes]:

    try:
        with tempfile.TemporaryDirectory() as tmp:
            doc_path = Path(tmp) / "input.doc"
            out_path = Path(tmp) / "input.docx"

            doc_path.write_bytes(file_bytes)

            subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "docx", str(doc_path), "--outdir", tmp],
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
                check=True
            )

            if out_path.exists():
                return out_path.read_bytes()

    except Exception:
        return None

    return None

def convert_image_to_pdf(file_bytes: bytes) -> Optional[bytes]:
    try:
        buf = io.BytesIO()
        Image.open(io.BytesIO(file_bytes)).convert("RGB").save(buf, format="PDF")
        return buf.getvalue()
    except Exception:
        return None

def convert_text_to_pdf(file_bytes: bytes) -> Optional[bytes]:
    try:
        text = file_bytes.decode("utf-8", errors="ignore").splitlines()
        buf = io.BytesIO()
        pdf = fitz.open()

        for i in range(0, len(text), TEXT_LINES_PER_PAGE):
            page = pdf.new_page()
            page.insert_text((72, 72), "\n".join(text[i:i + TEXT_LINES_PER_PAGE]))

        pdf.save(buf)
        return buf.getvalue()

    except Exception:
        return None


def convert_docx_to_pdf(file_bytes: bytes) -> Optional[bytes]:
    try:
        buf = io.BytesIO()
        pdf = fitz.open()
        doc = Document(io.BytesIO(file_bytes))

        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

        for i in range(0, len(paragraphs), PARAS_PER_PAGE):
            page = pdf.new_page()
            page.insert_text((72, 72), "\n".join(paragraphs[i:i + PARAS_PER_PAGE]))

        pdf.save(buf)
        return buf.getvalue()

    except Exception:
        return None

def convert_pptx_to_pdf(file_bytes: bytes) -> Optional[bytes]:
    try:
        buf = io.BytesIO()
        pdf = fitz.open()
        prs = Presentation(io.BytesIO(file_bytes))

        for slide in prs.slides:
            lines = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    lines.append(shape.text)

            page = pdf.new_page()
            page.insert_text((72, 72), "\n".join(lines))

        pdf.save(buf)
        return buf.getvalue()

    except Exception:
        return None
